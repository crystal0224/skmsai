"""FileAssembler — 콘텐츠 + 에셋 → 최종 파일 조립 서비스.

PR-046: PPTX/PDF/HTML/PNG/MP3 파일 조립.
PR-050: 테마 적용, 목차 슬라이드, 슬라이드 번호/푸터 삽입.

python-pptx로 강의자료 PPTX 생성.
Markdown → PDF, 이미지 세트, 오디오 합성 등.
"""
from __future__ import annotations

import logging
import os
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any

from scripts.lib.content_studio.models import (
    GeneratedAsset,
    GeneratedFile,
)
from scripts.lib.content_studio.pptx_templates import (
    PptxTheme,
    generate_footer,
    generate_toc_data,
    load_themes,
    select_layout,
    select_theme,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# 기업 테마 상수
# ---------------------------------------------------------------------------

SK_BLUE = "0052A2"
SK_RED = "E4002B"
SK_GRAY = "6C757D"
SK_LIGHT_GRAY = "F8F9FA"
WHITE = "FFFFFF"

# 슬라이드 크기 (16:9, EMU 단위)
SLIDE_WIDTH_EMU = 12192000  # 33.867 cm
SLIDE_HEIGHT_EMU = 6858000  # 19.05 cm

# 기본 폰트
DEFAULT_FONT = "맑은 고딕"
FALLBACK_FONT = "Arial"


# ---------------------------------------------------------------------------
# FileAssembler
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class AssemblerConfig:
    """FileAssembler 설정 — 불변.

    Attributes:
        output_dir: 출력 기본 디렉토리.
        theme_primary: 주 색상 (hex).
        theme_secondary: 보조 색상 (hex).
        font_name: 기본 폰트.
        aspect_ratio: 슬라이드 비율.
        style: 스타일 (professional, casual, academic).
        themes_yaml_path: pptx_themes.yaml 경로 (None이면 테마 미사용).
        company: 푸터 회사명.
    """

    output_dir: str = "output"
    theme_primary: str = SK_BLUE
    theme_secondary: str = SK_RED
    font_name: str = DEFAULT_FONT
    aspect_ratio: str = "16:9"
    style: str = "professional"
    themes_yaml_path: str | None = None
    company: str = "SK Group"


class FileAssembler:
    """콘텐츠 + 에셋 → 최종 파일 조립."""

    def __init__(self, config: AssemblerConfig | None = None) -> None:
        self._config = config or AssemblerConfig()
        self._theme = self._resolve_theme()
        self._ensure_output_dirs()

    def _resolve_theme(self) -> PptxTheme | None:
        """설정의 style + themes_yaml_path로 PptxTheme를 로드한다."""
        if not self._config.themes_yaml_path:
            return None
        try:
            themes = load_themes(self._config.themes_yaml_path)
            theme_name = select_theme(self._config.style)
            return themes.get(theme_name)
        except FileNotFoundError:
            logger.warning("테마 파일을 찾을 수 없습니다: %s", self._config.themes_yaml_path)
            return None

    @property
    def theme(self) -> PptxTheme | None:
        """현재 로드된 테마."""
        return self._theme

    def _ensure_output_dirs(self) -> None:
        """출력 디렉토리 자동 생성."""
        base = Path(self._config.output_dir)
        for subdir in (
            "lectures",
            "cardnews",
            "workshops",
            "visualizations",
            "audio",
            "quizzes",
            "assets",
        ):
            (base / subdir).mkdir(parents=True, exist_ok=True)

    def _sanitize_filename(self, topic: str) -> str:
        """파일명에 사용 불가한 문자 제거 + path traversal 방지."""
        # path traversal 시퀀스 제거
        sanitized = (
            topic.replace("..", "")
            .replace("/", "-")
            .replace("\\", "-")
            .replace(" ", "-")
        )
        sanitized = "".join(c for c in sanitized if c.isalnum() or c in "-_.")
        # 선행 점(.) 제거 (숨김파일 방지)
        sanitized = sanitized.lstrip(".")
        return sanitized[:50] if sanitized else "untitled"

    def _output_path(self, content_type: str, topic: str, ext: str) -> Path:
        """표준 출력 경로 생성: output/{type}/{topic}-{date}.{ext}.

        Raises:
            ValueError: 결과 경로가 output_dir 밖을 가리키는 경우.
        """
        subdir_map = {
            "lecture": "lectures",
            "card_news": "cardnews",
            "workshop": "workshops",
            "visualization": "visualizations",
            "audio": "audio",
            "quiz": "quizzes",
        }
        subdir = subdir_map.get(content_type, content_type)
        filename = f"{self._sanitize_filename(topic)}-{date.today().isoformat()}.{ext}"
        out = Path(self._config.output_dir) / subdir / filename
        # path traversal 방지: resolved 경로가 output_dir 내부인지 검증
        base_resolved = Path(self._config.output_dir).resolve()
        out_resolved = out.resolve()
        if not str(out_resolved).startswith(str(base_resolved)):
            raise ValueError(
                f"Path traversal 차단: {out} → {out_resolved} (base: {base_resolved})"
            )
        return out

    # ----- Lecture (PPTX) -----

    def assemble_lecture(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """강의자료 PPTX 생성.

        PR-050: 테마 색상, TOC 슬라이드, 슬라이드 번호/푸터 적용.

        Args:
            content: LectureContent (slides, citations).
            assets: 이미지/차트 에셋 목록.
            topic: 주제 (파일명에 사용).

        Returns:
            GeneratedFile (pptx).
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.warning("python-pptx 미설치. 텍스트 fallback 사용.")
            return self._assemble_lecture_fallback(content, topic)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH_EMU
        prs.slide_height = SLIDE_HEIGHT_EMU

        # 테마 색상 해석
        theme = self._theme
        title_color = self._parse_hex_color(theme.primary) if theme else None
        text_color = self._parse_hex_color(theme.text) if theme else None
        font_name = theme.font if theme else self._config.font_name

        # 에셋 인덱스 매핑
        asset_map: dict[int, GeneratedAsset] = {}
        if assets:
            for asset in assets:
                meta = dict(asset.metadata) if asset.metadata else {}
                slide_idx = meta.get("slide_index")
                if slide_idx is not None:
                    asset_map[int(slide_idx)] = asset

        slides = content.slides if hasattr(content, "slides") else ()

        # TOC 생성 (slides가 SlidePlan 속성을 가지면)
        toc_data = generate_toc_data(slides)

        # TOC 슬라이드 삽입 (TOC 항목이 2개 이상이면)
        if len(toc_data) >= 2:
            toc_layout = prs.slide_layouts[1]
            toc_slide = prs.slides.add_slide(toc_layout)
            if toc_slide.shapes.title:
                toc_slide.shapes.title.text = "목차"
                if title_color:
                    toc_slide.shapes.title.text_frame.paragraphs[
                        0
                    ].font.color.rgb = RGBColor(*title_color)
                toc_slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
            for shape in toc_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for i, entry in enumerate(toc_data):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"{entry['index']}. {entry['title']}"
                        p.font.size = Pt(16)
                        p.font.name = font_name
                        if text_color:
                            p.font.color.rgb = RGBColor(*text_color)
                    break

        # 콘텐츠 슬라이드
        footer_data = generate_footer(
            company=self._config.company,
        )
        slide_number = len(prs.slides) + 1  # TOC 다음부터

        for slide_content in slides:
            # 레이아웃 자동 선택
            layout_name = select_layout(slide_content)
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            # 제목
            if slide.shapes.title:
                slide.shapes.title.text = slide_content.title
                if title_color:
                    slide.shapes.title.text_frame.paragraphs[
                        0
                    ].font.color.rgb = RGBColor(*title_color)
                slide.shapes.title.text_frame.paragraphs[0].font.name = font_name

            # 본문
            body_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_placeholder = shape
                    break

            if body_placeholder is not None:
                tf = body_placeholder.text_frame
                tf.clear()
                # key points
                for i, point in enumerate(slide_content.key_points):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = point
                    p.font.size = Pt(16)
                    p.font.name = font_name
                    if text_color:
                        p.font.color.rgb = RGBColor(*text_color)

                # 본문 텍스트
                if slide_content.body_text:
                    p = tf.add_paragraph()
                    p.text = slide_content.body_text
                    p.font.size = Pt(14)
                    p.font.name = font_name
                    if text_color:
                        p.font.color.rgb = RGBColor(*text_color)

            # 이미지 삽입
            asset = asset_map.get(slide_content.index)
            if asset and os.path.exists(asset.file_path):
                try:
                    left = Inches(7.0)
                    top = Inches(1.5)
                    width = Inches(3.0)
                    slide.shapes.add_picture(asset.file_path, left, top, width=width)
                except Exception as e:
                    logger.warning(f"이미지 삽입 실패 (slide {slide_content.index}): {e}")

            # 발표자 노트
            if slide_content.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_content.speaker_notes

            # 푸터 (슬라이드 번호 + 회사명 + 날짜)
            if footer_data.get("show_page_number"):
                self._add_footer_textbox(
                    slide,
                    prs,
                    Pt,
                    Inches,
                    RGBColor,
                    slide_number,
                    footer_data,
                    font_name,
                    text_color,
                )
            slide_number += 1

        # 출처 슬라이드
        citations = content.citations if hasattr(content, "citations") else ()
        if citations:
            slide_layout = prs.slide_layouts[1]
            cite_slide = prs.slides.add_slide(slide_layout)
            if cite_slide.shapes.title:
                cite_slide.shapes.title.text = "출처"
                if title_color:
                    cite_slide.shapes.title.text_frame.paragraphs[
                        0
                    ].font.color.rgb = RGBColor(*title_color)
                cite_slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
            for shape in cite_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for i, cite in enumerate(citations):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"• {cite}"
                        p.font.size = Pt(12)
                        p.font.name = font_name
                    break

        # 저장
        out_path = self._output_path("lecture", topic, "pptx")
        prs.save(str(out_path))
        size = out_path.stat().st_size

        return GeneratedFile(
            file_type="pptx",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=size,
        )

    @staticmethod
    def _parse_hex_color(hex_str: str) -> tuple[int, int, int] | None:
        """'#RRGGBB' 또는 'RRGGBB' 문자열 → (R, G, B) 튜플."""
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return None
        try:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return None

    @staticmethod
    def _add_footer_textbox(
        slide: Any,
        prs: Any,
        Pt: Any,
        Inches: Any,
        RGBColor: Any,
        slide_number: int,
        footer_data: dict[str, Any],
        font_name: str,
        text_color: tuple[int, int, int] | None,
    ) -> None:
        """슬라이드 하단에 푸터 텍스트 박스를 추가한다."""
        footer_text = (
            f"{footer_data.get('company', '')}  |  "
            f"{footer_data.get('date', '')}  |  "
            f"{slide_number}"
        )
        left = Inches(0.5)
        top = prs.slide_height - Inches(0.5)
        width = Inches(9.0)
        height = Inches(0.3)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(8)
        p.font.name = font_name
        if text_color:
            p.font.color.rgb = RGBColor(*text_color)

    def _assemble_lecture_fallback(self, content: Any, topic: str) -> GeneratedFile:
        """python-pptx 미설치 시 텍스트 파일 fallback."""
        out_path = self._output_path("lecture", topic, "html")
        lines = [f"<html><body><h1>{topic}</h1>"]
        slides = content.slides if hasattr(content, "slides") else ()
        for sc in slides:
            lines.append(f"<h2>{sc.title}</h2>")
            for kp in sc.key_points:
                lines.append(f"<li>{kp}</li>")
            if sc.body_text:
                lines.append(f"<p>{sc.body_text}</p>")
        lines.append("</body></html>")
        text = "\n".join(lines)
        out_path.write_text(text, encoding="utf-8")
        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(text.encode("utf-8")),
        )

    # ----- Card News (PNG set) -----

    def assemble_card_news(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> tuple[GeneratedFile, ...]:
        """카드뉴스 PNG 세트 또는 HTML fallback."""
        results = []
        cards = content.cards if hasattr(content, "cards") else ()

        for card in cards:
            # 에셋 이미지가 있으면 사용
            asset = None
            if assets:
                for a in assets:
                    meta = dict(a.metadata) if a.metadata else {}
                    if meta.get("card_index") == card.index:
                        asset = a
                        break

            if asset and os.path.exists(asset.file_path):
                import shutil

                out_path = self._output_path(
                    "card_news", f"{topic}-{card.index}", "png"
                )
                shutil.copy2(asset.file_path, str(out_path))
                size = out_path.stat().st_size
                results.append(
                    GeneratedFile(
                        file_type="png",
                        file_path=str(out_path),
                        file_name=out_path.name,
                        size_bytes=size,
                    )
                )
            else:
                # HTML fallback
                out_path = self._output_path(
                    "card_news", f"{topic}-{card.index}", "html"
                )
                html = f"""<html><body style="width:1080px;height:1080px;background:#0052A2;color:white;padding:40px;">
<h1>{card.headline}</h1><p>{card.body}</p></body></html>"""
                out_path.write_text(html, encoding="utf-8")
                results.append(
                    GeneratedFile(
                        file_type="html",
                        file_path=str(out_path),
                        file_name=out_path.name,
                        size_bytes=len(html.encode("utf-8")),
                    )
                )

        return tuple(results)

    # ----- Workshop (PDF) -----

    def assemble_workshop(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """워크숍 PDF (Markdown → HTML fallback)."""
        phases = content.phases if hasattr(content, "phases") else ()
        lines = [f"# {topic} — 워크숍 진행자 가이드\n"]
        for phase in phases:
            lines.append(f"## {phase.title} ({phase.phase_type})")
            lines.append(phase.content_text)
            lines.append("")

        text = "\n".join(lines)
        out_path = self._output_path("workshop", topic, "html")
        out_path.write_text(text, encoding="utf-8")

        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(text.encode("utf-8")),
        )

    # ----- Audio (MP3 / script fallback) -----

    def assemble_audio(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """오디오 MP3 합성 또는 스크립트 텍스트 fallback."""
        # MP3 에셋이 있으면 합성 시도
        if assets:
            mp3_assets = [a for a in assets if a.asset_type == "audio"]
            if mp3_assets and all(os.path.exists(a.file_path) for a in mp3_assets):
                try:
                    return self._merge_audio(mp3_assets, topic)
                except Exception as e:
                    logger.warning(f"오디오 합성 실패: {e}")

        # 스크립트 텍스트 fallback
        sections = content.sections if hasattr(content, "sections") else ()
        lines = [f"# {topic} — 오디오 스크립트\n"]
        for sec in sections:
            lines.append(f"**[{sec.speaker}]**")
            lines.append(sec.text)
            lines.append("")

        text = "\n".join(lines)
        out_path = self._output_path("audio", topic, "html")
        out_path.write_text(text, encoding="utf-8")

        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(text.encode("utf-8")),
        )

    def _merge_audio(
        self, mp3_assets: list[GeneratedAsset], topic: str
    ) -> GeneratedFile:
        """MP3 파일 합성 (pydub 사용)."""
        from pydub import AudioSegment

        combined = AudioSegment.empty()
        for asset in mp3_assets:
            segment = AudioSegment.from_mp3(asset.file_path)
            combined += segment

        out_path = self._output_path("audio", topic, "mp3")
        combined.export(str(out_path), format="mp3")
        size = out_path.stat().st_size

        return GeneratedFile(
            file_type="mp3",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=size,
        )

    # ----- Visualization (SVG/PNG passthrough) -----

    def assemble_visualization(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """시각화 에셋 파일 복사."""
        if assets:
            chart_assets = [a for a in assets if a.asset_type == "chart"]
            if chart_assets:
                import shutil

                asset = chart_assets[0]
                ext = "svg" if asset.file_path.endswith(".svg") else "png"
                out_path = self._output_path("visualization", topic, ext)
                shutil.copy2(asset.file_path, str(out_path))
                return GeneratedFile(
                    file_type=ext,
                    file_path=str(out_path),
                    file_name=out_path.name,
                    size_bytes=out_path.stat().st_size,
                )

        # No asset fallback
        out_path = self._output_path("visualization", topic, "html")
        html = f"<html><body><h1>{topic}</h1><p>시각화 데이터 (차트 생성 불가)</p></body></html>"
        out_path.write_text(html, encoding="utf-8")
        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(html.encode("utf-8")),
        )
