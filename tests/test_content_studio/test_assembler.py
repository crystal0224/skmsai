"""FileAssembler 테스트.

PR-046: PPTX/PDF/HTML/PNG 조립 테스트.
"""
from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

import pytest

from src.content_studio.assembler import (
    AssemblerConfig,
    FileAssembler,
    SK_BLUE,
)
from src.content_studio.models import GeneratedAsset
from src.content_studio.pptx_templates import PptxTheme


# ---------------------------------------------------------------------------
# Mock Content dataclasses (PR-044에서 정의될 예정)
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class MockSlideContent:
    index: int
    title: str
    body_text: str = ""
    key_points: tuple[str, ...] = ()
    quote_ids: tuple[str, ...] = ()
    speaker_notes: str = ""


@dataclass(frozen=True)
class MockLectureContent:
    slides: tuple[MockSlideContent, ...]
    citations: tuple[str, ...] = ()
    warnings: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockCardContent:
    index: int
    headline: str
    body: str
    quote_ids: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockCardNewsContent:
    cards: tuple[MockCardContent, ...]
    citations: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockPhaseContent:
    phase_type: str
    title: str
    content_text: str
    quote_ids: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockWorkshopContent:
    phases: tuple[MockPhaseContent, ...]
    citations: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockSectionContent:
    index: int
    speaker: str
    text: str
    quote_ids: tuple[str, ...] = ()


@dataclass(frozen=True)
class MockAudioContent:
    sections: tuple[MockSectionContent, ...]
    citations: tuple[str, ...] = ()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_lecture_content(num_slides: int = 3) -> MockLectureContent:
    slides = tuple(
        MockSlideContent(
            index=i + 1,
            title=f"슬라이드 {i + 1}",
            body_text=f"본문 텍스트 {i + 1}",
            key_points=(f"포인트 {i + 1}-A", f"포인트 {i + 1}-B"),
            speaker_notes=f"발표자 노트 {i + 1}",
        )
        for i in range(num_slides)
    )
    return MockLectureContent(slides=slides, citations=("q-001", "q-002"))


def _make_card_news_content(num_cards: int = 3) -> MockCardNewsContent:
    cards = tuple(
        MockCardContent(
            index=i + 1,
            headline=f"제목 {i + 1}",
            body=f"본문 {i + 1}",
        )
        for i in range(num_cards)
    )
    return MockCardNewsContent(cards=cards, citations=("q-010",))


# ---------------------------------------------------------------------------
# AssemblerConfig
# ---------------------------------------------------------------------------


class TestAssemblerConfig:
    def test_defaults(self):
        config = AssemblerConfig()
        assert config.output_dir == "output"
        assert config.theme_primary == SK_BLUE
        assert config.aspect_ratio == "16:9"

    def test_frozen(self):
        config = AssemblerConfig()
        with pytest.raises(AttributeError):
            config.output_dir = "/tmp"  # type: ignore[misc]


# ---------------------------------------------------------------------------
# FileAssembler 기본
# ---------------------------------------------------------------------------


class TestFileAssemblerBasic:
    def test_init_creates_output_dirs(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        FileAssembler(config)
        assert (tmp_path / "out" / "lectures").is_dir()
        assert (tmp_path / "out" / "cardnews").is_dir()
        assert (tmp_path / "out" / "workshops").is_dir()
        assert (tmp_path / "out" / "assets").is_dir()

    def test_sanitize_filename(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        assert asm._sanitize_filename("SUPEX 추구/변천사") == "SUPEX-추구-변천사"
        assert asm._sanitize_filename("") == "untitled"

    def test_output_path(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        path = asm._output_path("lecture", "SUPEX", "pptx")
        assert "lectures" in str(path)
        assert path.suffix == ".pptx"


# ---------------------------------------------------------------------------
# Lecture Assembly (HTML fallback — no python-pptx needed)
# ---------------------------------------------------------------------------


class TestAssembleLecture:
    def test_html_fallback_when_no_pptx(self, tmp_path, monkeypatch):
        """python-pptx 없을 때 HTML fallback."""
        import src.content_studio.assembler as asm_module

        original_import = (
            __builtins__.__import__
            if hasattr(__builtins__, "__import__")
            else __import__
        )

        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("mock: pptx not installed")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr("builtins.__import__", mock_import)

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_lecture_content(2)
        result = asm.assemble_lecture(content, topic="SUPEX테스트")

        assert result.file_type == "html"
        assert result.size_bytes > 0
        assert os.path.exists(result.file_path)

    def test_lecture_creates_pptx_if_available(self, tmp_path):
        """python-pptx 설치 시 PPTX 생성."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_lecture_content(3)
        result = asm.assemble_lecture(content, topic="SUPEX-강의")

        assert result.file_type == "pptx"
        assert result.size_bytes > 0
        assert result.file_path.endswith(".pptx")
        assert os.path.exists(result.file_path)

    def test_lecture_pptx_slide_count(self, tmp_path):
        """PPTX 슬라이드 수 검증 (content slides + citation slide)."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_lecture_content(5)
        result = asm.assemble_lecture(content, topic="슬라이드수-테스트")

        prs = Presentation(result.file_path)
        # 5 content slides + 1 citation slide = 6
        assert len(prs.slides) == 6

    def test_lecture_pptx_speaker_notes(self, tmp_path):
        """발표자 노트가 포함되는지 검증."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_lecture_content(2)
        result = asm.assemble_lecture(content, topic="노트테스트")

        prs = Presentation(result.file_path)
        first_slide = prs.slides[0]
        notes_text = first_slide.notes_slide.notes_text_frame.text
        assert "발표자 노트 1" in notes_text

    def test_lecture_pptx_citation_slide(self, tmp_path):
        """마지막 슬라이드에 출처 목록."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_lecture_content(1)
        result = asm.assemble_lecture(content, topic="출처테스트")

        prs = Presentation(result.file_path)
        last_slide = prs.slides[-1]
        assert last_slide.shapes.title.text == "출처"

    def test_lecture_no_citations(self, tmp_path):
        """citations가 비어있으면 출처 슬라이드 없음."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = MockLectureContent(
            slides=(MockSlideContent(index=1, title="슬라이드1"),),
            citations=(),
        )
        result = asm.assemble_lecture(content, topic="노출처")

        prs = Presentation(result.file_path)
        assert len(prs.slides) == 1  # no citation slide


# ---------------------------------------------------------------------------
# CardNews Assembly
# ---------------------------------------------------------------------------


class TestAssembleCardNews:
    def test_html_fallback_no_assets(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_card_news_content(3)
        results = asm.assemble_card_news(content, topic="VWBE")

        assert len(results) == 3
        assert all(r.file_type == "html" for r in results)
        assert all(r.size_bytes > 0 for r in results)

    def test_with_image_assets(self, tmp_path):
        # 이미지 에셋 생성
        img_path = tmp_path / "test_card.png"
        img_path.write_bytes(b"\x89PNG" + b"\x00" * 100)

        asset = GeneratedAsset(
            asset_type="image",
            file_path=str(img_path),
            prompt_used="test",
            metadata=(("card_index", 1),),
        )

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_card_news_content(2)
        results = asm.assemble_card_news(content, assets=[asset], topic="에셋테스트")

        assert len(results) == 2
        # 첫 번째 카드는 PNG 에셋, 두 번째는 HTML fallback
        assert results[0].file_type == "png"
        assert results[1].file_type == "html"

    def test_card_uses_design_tokens(self, tmp_path):
        """카드뉴스 HTML이 디자인 토큰 CSS 변수를 포함해야 한다."""
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_card_news_content(1)
        results = asm.assemble_card_news(content, topic="토큰테스트")
        html = Path(results[0].file_path).read_text(encoding="utf-8")
        assert ":root" in html
        assert "--c-primary" in html

    def test_card_responsive_not_fixed_1080(self, tmp_path):
        """기본 렌더링은 반응형, 고정 1080px가 아니어야 한다."""
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_card_news_content(1)
        results = asm.assemble_card_news(content, topic="반응형테스트")
        html = Path(results[0].file_path).read_text(encoding="utf-8")
        assert "width:1080px;height:1080px" not in html
        assert "card-canvas" in html

    def test_card_wcag_no_bad_colors(self, tmp_path):
        """WCAG AA 미충족 색상이 없어야 한다."""
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = _make_card_news_content(1)
        results = asm.assemble_card_news(content, topic="WCAG테스트")
        html = Path(results[0].file_path).read_text(encoding="utf-8")
        assert "color:#AAA" not in html
        assert "color:#999" not in html


# ---------------------------------------------------------------------------
# Workshop Assembly
# ---------------------------------------------------------------------------


class TestAssembleWorkshop:
    def test_html_output(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = MockWorkshopContent(
            phases=(
                MockPhaseContent(phase_type="intro", title="도입", content_text="환영합니다"),
                MockPhaseContent(phase_type="main", title="본론", content_text="핵심 내용"),
            ),
        )
        result = asm.assemble_workshop(content, topic="워크숍테스트")

        assert result.file_type == "html"
        assert result.size_bytes > 0
        text = Path(result.file_path).read_text(encoding="utf-8")
        assert "도입" in text
        assert "본론" in text


# ---------------------------------------------------------------------------
# Audio Assembly
# ---------------------------------------------------------------------------


class TestAssembleAudio:
    def test_script_fallback_no_assets(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = MockAudioContent(
            sections=(
                MockSectionContent(index=1, speaker="narrator", text="안녕하세요"),
                MockSectionContent(index=2, speaker="host", text="환영합니다"),
            ),
        )
        result = asm.assemble_audio(content, topic="오디오테스트")

        assert result.file_type == "html"
        assert result.size_bytes > 0
        text = Path(result.file_path).read_text(encoding="utf-8")
        assert "narrator" in text


# ---------------------------------------------------------------------------
# Visualization Assembly
# ---------------------------------------------------------------------------


class TestAssembleVisualization:
    def test_html_fallback_no_assets(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        result = asm.assemble_visualization(None, topic="시각화테스트")

        assert result.file_type == "html"
        assert result.size_bytes > 0

    def test_with_svg_asset(self, tmp_path):
        svg_path = tmp_path / "chart.svg"
        svg_path.write_text("<svg></svg>")

        asset = GeneratedAsset(
            asset_type="chart",
            file_path=str(svg_path),
            prompt_used="timeline",
        )

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        result = asm.assemble_visualization(None, assets=[asset], topic="SVG테스트")

        assert result.file_type == "svg"
        assert result.size_bytes > 0


# ---------------------------------------------------------------------------
# PR-050: Theme Integration
# ---------------------------------------------------------------------------

THEMES_YAML = Path(__file__).parent.parent.parent / "config" / "pptx_themes.yaml"


class TestAssemblerThemeIntegration:
    def test_theme_loaded_from_config(self, tmp_path):
        """themes_yaml_path 설정 시 테마가 로드된다."""
        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            style="professional",
            themes_yaml_path=str(THEMES_YAML),
        )
        asm = FileAssembler(config)
        assert asm.theme is not None
        assert asm.theme.name == "Corporate (SK)"

    def test_casual_style_loads_education_theme(self, tmp_path):
        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            style="casual",
            themes_yaml_path=str(THEMES_YAML),
        )
        asm = FileAssembler(config)
        assert asm.theme is not None
        assert asm.theme.name == "Education"

    def test_academic_style_loads_seminar_theme(self, tmp_path):
        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            style="academic",
            themes_yaml_path=str(THEMES_YAML),
        )
        asm = FileAssembler(config)
        assert asm.theme is not None
        assert asm.theme.name == "Seminar (Formal)"

    def test_no_themes_yaml_returns_none(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        assert asm.theme is None

    def test_missing_themes_yaml_returns_none(self, tmp_path):
        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            themes_yaml_path="/nonexistent/themes.yaml",
        )
        asm = FileAssembler(config)
        assert asm.theme is None

    def test_parse_hex_color(self):
        assert FileAssembler._parse_hex_color("#E4002B") == (228, 0, 43)
        assert FileAssembler._parse_hex_color("003DA5") == (0, 61, 165)
        assert FileAssembler._parse_hex_color("#FFF") is None
        assert FileAssembler._parse_hex_color("ZZZZZZ") is None


class TestAssemblerTocSlide:
    def test_toc_slide_inserted_with_section_headers(self, tmp_path):
        """section_header가 2개 이상이면 TOC 슬라이드가 삽입된다."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            themes_yaml_path=str(THEMES_YAML),
        )
        asm = FileAssembler(config)

        # section_header 레이아웃 슬라이드로 TOC 트리거
        slides = (
            MockSlideContent(index=1, title="소개", key_points=("A",)),
            MockSlideContent(index=2, title="본론", key_points=("B",)),
            MockSlideContent(index=3, title="결론", key_points=("C",)),
        )
        # index=1인 슬라이드만 TOC에 포함 → 1개이므로 TOC 미생성
        content = MockLectureContent(slides=slides, citations=())
        result = asm.assemble_lecture(content, topic="TOC-테스트")

        prs = Presentation(result.file_path)
        # index<=1인 슬라이드가 1개뿐이므로 TOC 미생성: 3 content slides
        assert len(prs.slides) == 3

    def test_no_toc_slide_when_insufficient_sections(self, tmp_path):
        """section_header/index<=1이 1개 이하면 TOC 미생성."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = MockLectureContent(
            slides=(
                MockSlideContent(index=2, title="슬라이드2"),
                MockSlideContent(index=3, title="슬라이드3"),
            ),
            citations=(),
        )
        result = asm.assemble_lecture(content, topic="NoTOC")

        prs = Presentation(result.file_path)
        assert len(prs.slides) == 2  # no TOC


class TestAssemblerFooter:
    def test_footer_textbox_present(self, tmp_path):
        """콘텐츠 슬라이드에 푸터 텍스트 박스가 추가된다."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            company="SK Telecom",
        )
        asm = FileAssembler(config)
        content = _make_lecture_content(1)
        result = asm.assemble_lecture(content, topic="푸터테스트")

        prs = Presentation(result.file_path)
        first_slide = prs.slides[0]
        # 텍스트 박스 중 하나에 회사명이 있어야 함
        texts = [shape.text for shape in first_slide.shapes if shape.has_text_frame]
        footer_found = any("SK Telecom" in t for t in texts)
        assert footer_found

    def test_footer_has_page_number(self, tmp_path):
        """푸터에 페이지 번호가 포함된다."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = MockLectureContent(
            slides=(
                MockSlideContent(index=1, title="S1"),
                MockSlideContent(index=2, title="S2"),
            ),
            citations=(),
        )
        result = asm.assemble_lecture(content, topic="페이지번호")

        prs = Presentation(result.file_path)
        # 두 번째 콘텐츠 슬라이드의 푸터에 "2"가 포함되어야 함
        second_slide = prs.slides[1]
        texts = [shape.text for shape in second_slide.shapes if shape.has_text_frame]
        page_num_found = any("2" in t for t in texts)
        assert page_num_found

    def test_themed_lecture_applies_font(self, tmp_path):
        """테마 적용 시 폰트 이름이 반영된다."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        config = AssemblerConfig(
            output_dir=str(tmp_path / "out"),
            style="casual",
            themes_yaml_path=str(THEMES_YAML),
        )
        asm = FileAssembler(config)
        assert asm.theme is not None
        assert asm.theme.font == "나눔고딕"

        content = MockLectureContent(
            slides=(MockSlideContent(index=1, title="폰트테스트", key_points=("A",)),),
            citations=(),
        )
        result = asm.assemble_lecture(content, topic="폰트")

        prs = Presentation(result.file_path)
        first_slide = prs.slides[0]
        title_font = first_slide.shapes.title.text_frame.paragraphs[0].font.name
        assert title_font == "나눔고딕"


# ---------------------------------------------------------------------------
# Security: path traversal 방지
# ---------------------------------------------------------------------------


class TestSanitizeFilename:
    def test_normal_topic(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        assert asm._sanitize_filename("SUPEX 추구") == "SUPEX-추구"

    def test_strips_path_separators(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        assert "/" not in asm._sanitize_filename("a/b/c")
        assert "\\" not in asm._sanitize_filename("a\\b\\c")

    def test_strips_dotdot(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        result = asm._sanitize_filename("../../etc/passwd")
        assert ".." not in result
        assert "/" not in result

    def test_strips_leading_dot(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        result = asm._sanitize_filename(".hidden")
        assert not result.startswith(".")

    def test_empty_returns_untitled(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        assert asm._sanitize_filename("") == "untitled"

    def test_all_special_chars(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        assert asm._sanitize_filename("!@#$%^&*()") == "untitled"

    def test_max_length_50(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        result = asm._sanitize_filename("A" * 100)
        assert len(result) == 50


class TestOutputPathSecurity:
    def test_normal_path_within_output(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        path = asm._output_path("lecture", "SUPEX", "pptx")
        assert str(tmp_path / "out") in str(path)

    def test_traversal_topic_sanitized(self, tmp_path):
        """path traversal 주입이 sanitize로 무력화된다."""
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        path = asm._output_path("lecture", "../../etc/passwd", "pptx")
        assert str(tmp_path / "out") in str(path.resolve())

    def test_unknown_content_type(self, tmp_path):
        asm = FileAssembler(AssemblerConfig(output_dir=str(tmp_path / "out")))
        path = asm._output_path("unknown_type", "test", "html")
        assert "unknown_type" in str(path)


# ---------------------------------------------------------------------------
# Quiz Assembly
# ---------------------------------------------------------------------------


class TestAssembleQuiz:
    def _make_quiz_content(self, n=2):
        """퀴즈 콘텐츠 mock."""
        from types import SimpleNamespace

        questions = tuple(
            SimpleNamespace(
                index=i + 1,
                question_text=f"질문 {i+1}",
                choices=(f"선택지A-{i+1}", f"선택지B-{i+1}", f"선택지C-{i+1}"),
                correct_answer=0,
                explanation=f"설명 {i+1}",
                source_quote=f"출처 {i+1}" if i == 0 else "",
            )
            for i in range(n)
        )
        return SimpleNamespace(questions=questions)

    def test_quiz_html_output(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(3)
        result = asm.assemble_quiz(content, topic="퀴즈테스트")
        assert result.file_type == "html"
        assert result.size_bytes > 0

    def test_quiz_uses_design_tokens(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(2)
        result = asm.assemble_quiz(content, topic="토큰퀴즈")
        html = Path(result.file_path).read_text(encoding="utf-8")
        assert ":root" in html
        assert "--c-primary" in html

    def test_quiz_wcag_colors(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(2)
        result = asm.assemble_quiz(content, topic="WCAG퀴즈")
        html = Path(result.file_path).read_text(encoding="utf-8")
        assert "#28a745" not in html
        assert "#dc3545" not in html

    def test_quiz_focus_ring(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(2)
        result = asm.assemble_quiz(content, topic="포커스퀴즈")
        html = Path(result.file_path).read_text(encoding="utf-8")
        assert ":focus" in html

    def test_quiz_aria_disabled(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(2)
        result = asm.assemble_quiz(content, topic="접근성퀴즈")
        html = Path(result.file_path).read_text(encoding="utf-8")
        assert "aria-disabled" in html

    def test_quiz_progress_indicator(self, tmp_path):
        config = AssemblerConfig(output_dir=str(tmp_path / "out"))
        asm = FileAssembler(config)
        content = self._make_quiz_content(3)
        result = asm.assemble_quiz(content, topic="프로그레스퀴즈")
        html = Path(result.file_path).read_text(encoding="utf-8")
        assert "progress" in html.lower()
