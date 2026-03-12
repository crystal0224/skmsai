"""FileAssembler — 콘텐츠 + 에셋 → 최종 파일 조립 서비스.

PR-046: PPTX/PDF/HTML/PNG/MP3 파일 조립.
PR-050: 테마 적용, 목차 슬라이드, 슬라이드 번호/푸터 삽입.

python-pptx로 강의자료 PPTX 생성.
Markdown → PDF, 이미지 세트, 오디오 합성 등.
"""
from __future__ import annotations

import logging
import os
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any

from src.content_studio.models import (
    GeneratedAsset,
    GeneratedFile,
)
from src.content_studio.pptx_templates import (
    PptxTheme,
    generate_footer,
    generate_toc_data,
    load_themes,
    select_layout,
    select_theme,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# 기업 테마 상수
# ---------------------------------------------------------------------------

SK_BLUE = "0052A2"
SK_RED = "E4002B"
SK_GRAY = "6C757D"
SK_LIGHT_GRAY = "F8F9FA"
WHITE = "FFFFFF"

# ---------------------------------------------------------------------------
# 디자인 시스템 규격 (Typography Hierarchy)
# ---------------------------------------------------------------------------

FONT_SIZE_TITLE = 24
FONT_SIZE_GOVERNING = 18
FONT_SIZE_KEY_POINT = 16
FONT_SIZE_BODY = 14
FONT_SIZE_CITATION = 11
FONT_SIZE_FOOTER = 8

# 슬라이드 크기 (16:9, EMU 단위)
SLIDE_WIDTH_EMU = 12192000  # 33.867 cm
SLIDE_HEIGHT_EMU = 6858000  # 19.05 cm

# 기본 폰트 (Pretendard 선호, 서버는 NanumGothic)
DEFAULT_FONT = "Pretendard"
SERVER_FONT = "NanumGothic"
FALLBACK_FONT = "맑은 고딕"

# Pretendard 웹폰트 (HTML용 CDN)
PRETENDARD_CSS = '<link rel="stylesheet" href="https://cdn.jsdelivr.net/gh/orioncactus/pretendard@v1.3.9/dist/web/static/pretendard.min.css" />'


# ---------------------------------------------------------------------------
# FileAssembler
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class AssemblerConfig:
    """FileAssembler 설정 — 불변.

    Attributes:
        output_dir: 출력 기본 디렉토리.
        theme_primary: 주 색상 (hex).
        theme_secondary: 보조 색상 (hex).
        font_name: 기본 폰트.
        aspect_ratio: 슬라이드 비율.
        style: 스타일 (professional, casual, academic).
        themes_yaml_path: pptx_themes.yaml 경로 (None이면 테마 미사용).
        company: 푸터 회사명.
    """

    output_dir: str = "output"
    theme_primary: str = SK_BLUE
    theme_secondary: str = SK_RED
    font_name: str = DEFAULT_FONT
    aspect_ratio: str = "16:9"
    style: str = "professional"
    themes_yaml_path: str | None = None
    company: str = "SK Group"


class FileAssembler:
    """콘텐츠 + 에셋 → 최종 파일 조립."""

    def __init__(self, config: AssemblerConfig | None = None) -> None:
        self._config = config or AssemblerConfig()
        self._theme = self._resolve_theme()
        self._ensure_output_dirs()

    def _resolve_theme(self) -> PptxTheme | None:
        """설정의 style + themes_yaml_path로 PptxTheme를 로드한다."""
        if not self._config.themes_yaml_path:
            return None
        try:
            themes = load_themes(self._config.themes_yaml_path)
            theme_name = select_theme(self._config.style)
            return themes.get(theme_name)
        except FileNotFoundError:
            logger.warning("테마 파일을 찾을 수 없습니다: %s", self._config.themes_yaml_path)
            return None

    @property
    def theme(self) -> PptxTheme | None:
        """현재 로드된 테마."""
        return self._theme

    def _ensure_output_dirs(self) -> None:
        """출력 디렉토리 자동 생성."""
        base = Path(self._config.output_dir)
        for subdir in (
            "lectures",
            "cardnews",
            "workshops",
            "visualizations",
            "audio",
            "quizzes",
            "assets",
        ):
            (base / subdir).mkdir(parents=True, exist_ok=True)

    def _sanitize_filename(self, topic: str) -> str:
        """파일명에 사용 불가한 문자 제거 + path traversal 방지."""
        # path traversal 시퀀스 제거
        sanitized = (
            topic.replace("..", "")
            .replace("/", "-")
            .replace("\\", "-")
            .replace(" ", "-")
        )
        sanitized = "".join(c for c in sanitized if c.isalnum() or c in "-_.")
        # 선행 점(.) 제거 (숨김파일 방지)
        sanitized = sanitized.lstrip(".")
        return sanitized[:50] if sanitized else "untitled"

    def _output_path(self, content_type: str, topic: str, ext: str) -> Path:
        """표준 출력 경로 생성: output/{type}/{topic}-{date}.{ext}.

        Raises:
            ValueError: 결과 경로가 output_dir 밖을 가리키는 경우.
        """
        subdir_map = {
            "lecture": "lectures",
            "card_news": "cardnews",
            "workshop": "workshops",
            "visualization": "visualizations",
            "audio": "audio",
            "quiz": "quizzes",
        }
        subdir = subdir_map.get(content_type, content_type)
        filename = f"{self._sanitize_filename(topic)}-{date.today().isoformat()}.{ext}"
        out = Path(self._config.output_dir) / subdir / filename
        # path traversal 방지: resolved 경로가 output_dir 내부인지 검증
        base_resolved = Path(self._config.output_dir).resolve()
        out_resolved = out.resolve()
        if not str(out_resolved).startswith(str(base_resolved)):
            raise ValueError(
                f"Path traversal 차단: {out} → {out_resolved} (base: {base_resolved})"
            )
        return out

    # ----- Lecture (PPTX) -----

    def assemble_lecture(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """강의자료 PPTX 생성.

        PR-050: 테마 색상, TOC 슬라이드, 슬라이드 번호/푸터 적용.

        Args:
            content: LectureContent (slides, citations).
            assets: 이미지/차트 에셋 목록.
            topic: 주제 (파일명에 사용).

        Returns:
            GeneratedFile (pptx).
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN

            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.warning("python-pptx 미설치. 텍스트 fallback 사용.")
            return self._assemble_lecture_fallback(content, topic)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH_EMU
        prs.slide_height = SLIDE_HEIGHT_EMU

        # 테마 색상 해석
        theme = self._theme
        title_color = self._parse_hex_color(theme.primary) if theme else None
        text_color = self._parse_hex_color(theme.text) if theme else None
        font_name = theme.font if theme else self._config.font_name
        # [폰트 고도화] 리눅스 서버(Docker) 환경이면 설치된 NanumGothic 강제 적용
        if os.name != 'nt':
            font_name = "NanumGothic"

        # 에셋 인덱스 매핑
        asset_map: dict[int, GeneratedAsset] = {}
        if assets:
            for asset in assets:
                meta = dict(asset.metadata) if asset.metadata else {}
                slide_idx = meta.get("slide_index")
                if slide_idx is not None:
                    asset_map[int(slide_idx)] = asset

        slides = content.slides if hasattr(content, "slides") else ()

        # TOC 생성 (slides가 SlidePlan 속성을 가지면)
        toc_data = generate_toc_data(slides)

        # TOC 슬라이드 삽입 (TOC 항목이 2개 이상이면)
        if len(toc_data) >= 2:
            toc_layout = prs.slide_layouts[1]
            toc_slide = prs.slides.add_slide(toc_layout)
            if toc_slide.shapes.title:
                toc_slide.shapes.title.text = "목차"
                if title_color:
                    toc_slide.shapes.title.text_frame.paragraphs[
                        0
                    ].font.color.rgb = RGBColor(*title_color)
                toc_slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
            for shape in toc_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for i, entry in enumerate(toc_data):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"{entry['index']}. {entry['title']}"
                        p.font.size = Pt(16)
                        p.font.name = font_name
                        if text_color:
                            p.font.color.rgb = RGBColor(*text_color)
                    break

        # 콘텐츠 슬라이드
        footer_data = generate_footer(
            company=self._config.company,
        )
        slide_number = len(prs.slides) + 1  # TOC 다음부터

        for slide_content in slides:
            # 레이아웃 및 도식화 판단 (PR-061: Semantic Diagrams)
            layout_name = select_layout(slide_content)
            hint = (getattr(slide_content, "design_hint", "") or "").lower()
            
            # 도식화 트리거 조건
            is_process = any(x in hint for x in ["프로세스", "단계", "process", "flow"])
            is_matrix = any(x in hint for x in ["매트릭스", "4분면", "matrix", "quadrant"])

            layout_idx = get_layout_index(layout_name)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # 제목 및 거버닝 메시지 고정 렌더링 (PR-057: Fixed Header/Governing Message)
            self._add_fixed_title_and_governing(
                slide, 
                prs,
                slide_content.title,
                getattr(slide_content, "governing_message", ""),
                font_name,
                title_color,
                text_color,
                Pt,
                Inches,
                RGBColor
            )

            key_points = slide_content.key_points if hasattr(slide_content, "key_points") else ()
            body_text = getattr(slide_content, "body_text", "")

            # [Task 2] 도식화 엔진 실행 또는 일반 텍스트 배치
            if is_process and 3 <= len(key_points) <= 5:
                self._draw_process_diagram(slide, key_points, font_name, Pt, Inches, RGBColor)
            elif is_matrix and len(key_points) == 4:
                self._draw_matrix_diagram(slide, key_points, font_name, Pt, Inches, RGBColor)
            else:
                # 기존 텍스트 배치 로직
                placeholders = sorted(
                    [s for s in slide.placeholders if s.placeholder_format.idx in (1, 2)],
                    key=lambda x: x.placeholder_format.idx
                )
                if layout_name in ("comparison", "two_content") and len(placeholders) >= 2:
                    mid = (len(key_points) + 1) // 2
                    self._fill_placeholder(placeholders[0], key_points[:mid], "", font_name, text_color, Pt, RGBColor)
                    self._fill_placeholder(placeholders[1], key_points[mid:], "", font_name, text_color, Pt, RGBColor)
                elif placeholders:
                    self._fill_placeholder(placeholders[0], key_points, body_text, font_name, text_color, Pt, RGBColor)
...

    def _add_fixed_title_and_governing(
        self,
        slide: Any,
        prs: Any,
        title: str,
        governing: str,
        font_name: str,
        title_color: tuple[int, int, int] | None,
        text_color: tuple[int, int, int] | None,
        Pt: Any,
        Inches: Any,
        RGBColor: Any
    ) -> None:
        """모든 슬라이드에서 제목과 거버닝 메시지를 장식 없이 고정된 위치와 크기로 렌더링한다."""
        # 1. 제목 (Title) - 24pt Bold (Blue 강조)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), prs.slide_width - Inches(1.0), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(FONT_SIZE_TITLE)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0, 82, 162) # #0052A2 (SK Blue)

        # 2. 거버닝 메시지 (Governing Message) - 18pt (Dark Gray)
        if governing:
            # 장식 요소 제거: 시각적 피로도를 줄이기 위해 여백으로만 구분
            gov_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), prs.slide_width - Inches(1.0), Inches(0.5))
            tf_gov = gov_box.text_frame
            tf_gov.word_wrap = True
            p_gov = tf_gov.paragraphs[0]
            p_gov.text = governing
            p_gov.font.size = Pt(FONT_SIZE_GOVERNING)
            p_gov.font.name = font_name
            p_gov.font.color.rgb = RGBColor(26, 26, 46) # #1A1A2E (Dark Gray)

            if layout_name in ("comparison", "two_content") and len(placeholders) >= 2:
                # 데이터를 반으로 나눔
                mid = (len(key_points) + 1) // 2
                left_points = key_points[:mid]
                right_points = key_points[mid:]
                
                # 왼쪽 placeholder
                self._fill_placeholder(placeholders[0], left_points, "", font_name, text_color, Pt, RGBColor)
                # 오른쪽 placeholder
                self._fill_placeholder(placeholders[1], right_points, "", font_name, text_color, Pt, RGBColor)
            elif placeholders:
                # 기본 단일 placeholder
                self._fill_placeholder(placeholders[0], key_points, body_text, font_name, text_color, Pt, RGBColor)

            # [Task 4] 핵심 개념 아이콘 추가 (PR-062: Icon System)
            self._add_concept_icon(slide, slide_content.title, Inches, RGBColor)

            # 이미지 삽입
...
    def _add_concept_icon(self, slide: Any, title: str, Inches: Any, RGBColor: Any) -> None:
        """제목 키워드를 분석하여 상징적인 아이콘(도형)을 추가한다."""
        from pptx.enum.shapes import MSO_SHAPE
        
        # 키워드별 도형 매핑 (미니멀 전문 디자인용)
        icon_map = {
            "행복": MSO_SHAPE.HEART,
            "인간": MSO_SHAPE.SMILEY_FACE,
            "VWBE": MSO_SHAPE.GEAR,
            "SUPEX": MSO_SHAPE.STAR_5_POINT,
            "이해관계자": MSO_SHAPE.FLOWCHART_CONNECTOR,
            "경영": MSO_SHAPE.CHART_UP,
            "원칙": MSO_SHAPE.SEAL_4,
            "조직": MSO_SHAPE.HEXAGON,
            "구성원": MSO_SHAPE.OVAL
        }
        
        target_shape = None
        for kw, shape_type in icon_map.items():
            if kw in title:
                target_shape = shape_type
                break
        
        if target_shape:
            # 우측 상단에 작고 세련되게 배치 (Inches(9.2), Inches(0.4))
            icon_size = Inches(0.3)
            shape = slide.shapes.add_shape(target_shape, slide.parent.slide_width - Inches(0.8), Inches(0.45), icon_size, icon_size)
            shape.fill.solid()
            # 은은한 라이트 블루 계열로 브랜드 강조
            shape.fill.fore_color.rgb = RGBColor(0, 82, 162) # SK Blue
            shape.line.visible = False
            shape.shadow.inherit = False # 그림자 제거로 미니멀 유지
            asset = asset_map.get(slide_content.index)
            if asset and os.path.exists(asset.file_path):
                try:
                    # 레이아웃에 맞춰 이미지 위치 조정
                    left = Inches(7.0) if layout_name == "title_content_image" else Inches(1.0)
                    top = Inches(1.5)
                    width = Inches(3.0)
                    slide.shapes.add_picture(asset.file_path, left, top, width=width)
                except Exception as e:
                    logger.warning(f"이미지 삽입 실패 (slide {slide_content.index}): {e}")

            # 발표자 노트 및 상세 출처 (PR-063: Deep Grounding)
            speaker_notes = getattr(slide_content, "speaker_notes", "")
            source_info = ""
            if hasattr(slide_content, "source_details") and slide_content.source_details:
                source_info = "\n\n" + "-"*30 + "\n📜 근거 원문 상세:\n"
                for detail in slide_content.source_details:
                    source_info += f"• [{detail['edition']}] {detail['page']}p: {detail['text']}\n"
            
            full_notes = speaker_notes + source_info
            if full_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = full_notes

            # 푸터 (슬라이드 번호 + 회사명 + 날짜)
            if footer_data.get("show_page_number"):
                self._add_footer_textbox(
                    slide,
                    prs,
                    Pt,
                    Inches,
                    RGBColor,
                    slide_number,
                    footer_data,
                    font_name,
                    text_color,
                )
            slide_number += 1

        # 출처 슬라이드
        citations = content.citations if hasattr(content, "citations") else ()
        if citations:
            slide_layout = prs.slide_layouts[1]
            cite_slide = prs.slides.add_slide(slide_layout)
            if cite_slide.shapes.title:
                cite_slide.shapes.title.text = "출처"
                if title_color:
                    cite_slide.shapes.title.text_frame.paragraphs[
                        0
                    ].font.color.rgb = RGBColor(*title_color)
                cite_slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
            for shape in cite_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for i, cite in enumerate(citations):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"• {cite}"
                        p.font.size = Pt(12)
                        p.font.name = font_name
                    break

        # 저장
        out_path = self._output_path("lecture", topic, "pptx")
        prs.save(str(out_path))
        size = out_path.stat().st_size

        return GeneratedFile(
            file_type="pptx",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=size,
        )

    @staticmethod
    def _parse_hex_color(hex_str: str) -> tuple[int, int, int] | None:
        """'#RRGGBB' 또는 'RRGGBB' 문자열 → (R, G, B) 튜플."""
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return None
        try:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return None

    @staticmethod
    def _add_footer_textbox(
        slide: Any,
        prs: Any,
        Pt: Any,
        Inches: Any,
        RGBColor: Any,
        slide_number: int,
        footer_data: dict[str, Any],
        font_name: str,
        text_color: tuple[int, int, int] | None,
    ) -> None:
        """슬라이드 하단에 푸터 텍스트 박스를 추가한다."""
        footer_text = (
            f"{footer_data.get('company', '')}  |  "
            f"{footer_data.get('date', '')}  |  "
            f"{slide_number}"
        )
        left = Inches(0.5)
        top = prs.slide_height - Inches(0.5)
        width = Inches(9.0)
        height = Inches(0.3)
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(8)
        p.font.name = font_name
        if text_color:
            p.font.color.rgb = RGBColor(*text_color)

    def _assemble_lecture_fallback(self, content: Any, topic: str) -> GeneratedFile:
        """python-pptx 미설치 시 텍스트 파일 fallback."""
        out_path = self._output_path("lecture", topic, "html")
        lines = [f"<html><body><h1>{topic}</h1>"]
        slides = content.slides if hasattr(content, "slides") else ()
        for sc in slides:
            lines.append(f"<h2>{sc.title}</h2>")
            for kp in sc.key_points:
                lines.append(f"<li>{kp}</li>")
            if sc.body_text:
                lines.append(f"<p>{sc.body_text}</p>")
        lines.append("</body></html>")
        text = "\n".join(lines)
        out_path.write_text(text, encoding="utf-8")
        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(text.encode("utf-8")),
        )

    def _fill_placeholder(self, shape: Any, key_points: tuple[str, ...], body_text: str, font_name: str, text_color: tuple[int, int, int] | None, Pt: Any, RGBColor: Any) -> None:
        """Placeholder 텍스트 프레임을 장식 없이 체계적인 위계로 채운다."""
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # 1. 핵심 포인트 (Key Points) - 16pt Bold (Blue 강조)
        for i, point in enumerate(key_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(FONT_SIZE_KEY_POINT)
            p.font.bold = True
            p.font.name = font_name
            # 포인트 텍스트도 블루 계열로 통일하여 시각적 흐름 유지
            p.font.color.rgb = RGBColor(0, 82, 162) # #0052A2
            p.space_after = Pt(8)

        # 2. 본문 텍스트 (Body Text) - 14pt (Dark Gray)
        if body_text:
            p = tf.add_paragraph()
            p.text = body_text
            p.font.size = Pt(FONT_SIZE_BODY)
            p.font.bold = False
            p.font.name = font_name
            p.font.color.rgb = RGBColor(26, 26, 46) # #1A1A2E
            p.space_before = Pt(12) # 정보 덩어리 사이의 충분한 여백

    def _draw_process_diagram(self, slide: Any, points: tuple[str, ...], font_name: str, Pt: Any, Inches: Any, RGBColor: Any) -> None:
        """단계별 화살표 프로세스 도식을 그린다."""
        n = min(len(points), 5)
        margin_x = Inches(0.5)
        avail_width = slide.parent.slide_width - (margin_x * 2)
        box_width = (avail_width / n) - Inches(0.2)
        box_height = Inches(1.8)
        top = Inches(2.2)

        for i in range(n):
            left = margin_x + (i * (box_width + Inches(0.2)))
            # 화살표 오각형 (Pentagon)
            shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, box_width, box_height)
            shape.fill.solid()
            # 그라데이션 블루: 1단계(연함) -> n단계(진함)
            blue_val = 200 - (i * 30)
            shape.fill.fore_color.rgb = RGBColor(0, 82, max(blue_val, 100))
            shape.line.visible = False
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = f"STEP {i+1}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = points[i]
            p2.font.size = Pt(14)
            p2.font.bold = False
            p2.font.color.rgb = RGBColor(255, 255, 255)

    def _draw_matrix_diagram(self, slide: Any, points: tuple[str, ...], font_name: str, Pt: Any, Inches: Any, RGBColor: Any) -> None:
        """2x2 매트릭스 도식을 그린다."""
        margin_left = (slide.parent.slide_width - Inches(6.0)) / 2
        top = Inches(2.0)
        box_size = Inches(2.8)
        gap = Inches(0.15)

        coords = [
            (margin_left, top), (margin_left + box_size + gap, top),
            (margin_left, top + box_size + gap), (margin_left + box_size + gap, top + box_size + gap)
        ]
        
        # 4개 영역 색상 분화
        colors = [(0, 82, 162), (40, 110, 180), (80, 140, 200), (120, 170, 220)]

        for i in range(min(len(points), 4)):
            left, t = coords[i]
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, t, box_size, box_size)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors[i])
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = points[i]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.name = font_name
            p.font.color.rgb = RGBColor(255, 255, 255)

    # ----- Card News (PNG set) -----

    def assemble_card_news(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> tuple[GeneratedFile, ...]:
        """카드뉴스 PNG 세트 또는 HTML fallback.

        PR-055: design_hint 기반 스타일 차별화.
        """
        results = []
        cards = content.cards if hasattr(content, "cards") else ()

        for card in cards:
            # 에셋 이미지가 있으면 사용
            asset = None
            if assets:
                for a in assets:
                    meta = dict(a.metadata) if a.metadata else {}
                    if meta.get("card_index") == card.index:
                        asset = a
                        break

            if asset and os.path.exists(asset.file_path):
                import shutil

                out_path = self._output_path(
                    "card_news", f"{topic}-{card.index}", "png"
                )
                shutil.copy2(asset.file_path, str(out_path))
                size = out_path.stat().st_size
                results.append(
                    GeneratedFile(
                        file_type="png",
                        file_path=str(out_path),
                        file_name=out_path.name,
                        size_bytes=size,
                    )
                )
            else:
                # HTML fallback (미니멀 디자인 시스템 적용)
                out_path = self._output_path(
                    "card_news", f"{topic}-{card.index}", "html"
                )
                body_html = card.body.replace("\n", "<br>")
                
                # 디자인 시스템 고정: 화이트 배경 + 블루 강조
                bg_style = "#FFFFFF"
                accent_color = "#0052A2" # SK Blue
                text_color = "#1A1A2E"   # Dark Gray
                border_color = "rgba(0, 82, 162, 0.2)"

                # 인용 블록
                quote_html = ""
                source_quote = getattr(card, "source_quote", "")
                source_edition = getattr(card, "source_edition", "")
                if source_quote:
                    quote_html = (
                        f'<blockquote style="border-left:4px solid {accent_color};'
                        f"padding-left:20px;margin-top:30px;font-style:italic;"
                        f'font-size:20px;color:#4A4A4A;line-height:1.6;background:#F8F9FA;padding-top:15px;padding-bottom:15px;">'
                        f"「{source_quote}」"
                        f'<span style="display:block;font-size:14px;margin-top:10px;color:{accent_color};font-weight:600;">'
                        f"— {source_edition}</span>"
                        f"</blockquote>"
                    )
                
                card_num = f'<div style="position:absolute;top:40px;right:60px;font-size:16px;color:#999;font-weight:600;">{card.index} / {len(cards)}</div>'
                html = f"""<!DOCTYPE html>
<html lang="ko">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">{PRETENDARD_CSS}</head>
<body style="margin:0;width:1080px;height:1080px;background:{bg_style};color:{text_color};font-family:'Pretendard','Apple SD Gothic Neo',sans-serif;display:flex;flex-direction:column;justify-content:center;padding:100px;box-sizing:border-box;position:relative;border:20px solid {accent_color};">
{card_num}
<div style="width:80px;height:4px;background:{accent_color};margin-bottom:30px;"></div>
<h1 style="font-size:52px;margin:0 0 30px;line-height:1.2;font-weight:800;color:{accent_color}; letter-spacing:-1px;">{card.headline}</h1>
<p style="font-size:28px;line-height:1.7;margin:0 0 20px;color:#333;font-weight:400;">{body_html}</p>
{quote_html}
<div style="position:absolute;bottom:60px;left:100px;font-size:14px;color:#AAA;letter-spacing:2px;font-weight:600;">SKMS CONTENT STUDIO</div>
</body></html>"""
                out_path.write_text(html, encoding="utf-8")
                results.append(
                    GeneratedFile(
                        file_type="html",
                        file_path=str(out_path),
                        file_name=out_path.name,
                        size_bytes=len(html.encode("utf-8")),
                    )
                )

        return tuple(results)

    def assemble_workshop(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """워크숍 가이드 HTML (미니멀 디자인 시스템)."""
        phases = content.phases if hasattr(content, "phases") else ()
        
        phase_html = ""
        for phase in phases:
            phase_html += f"""
            <div style="margin-bottom:40px; padding:30px; border-left:4px solid #0052A2; background:#F8F9FA;">
                <div style="color:#0052A2; font-weight:800; font-size:14px; margin-bottom:10px; letter-spacing:1px;">{phase.phase_type.upper()} PHASE</div>
                <h2 style="margin:0 0 15px; font-size:24px; color:#1A1A2E;">{phase.title}</h2>
                <div style="font-size:16px; line-height:1.8; color:#444;">{phase.content_text.replace(chr(10), '<br>')}</div>
            </div>
            """

        html = f"""<!DOCTYPE html>
<html lang="ko">
<head><meta charset="utf-8">
<style>
    body{{font-family:'Pretendard','Apple SD Gothic Neo',sans-serif; background:#FFFFFF; color:#1A1A2E; margin:0; padding:60px; line-height:1.6;}}
    header{{border-bottom:2px solid #EEE; margin-bottom:50px; padding-bottom:30px;}}
    h1{{color:#0052A2; font-size:36px; margin:0 0 10px; font-weight:800;}}
    .meta{{color:#666; font-size:16px;}}
    footer{{margin-top:100px; color:#AAA; font-size:12px; letter-spacing:1px; border-top:1px solid #EEE; padding-top:20px;}}
</style>
</head>
<body>
    <header>
        <h1>{topic}</h1>
        <div class="meta">SKMS 워크숍 진행자 가이드</div>
    </header>
    {phase_html}
    <footer>SKMS CONTENT STUDIO</footer>
</body></html>"""

        out_path = self._output_path("workshop", topic, "html")
        out_path.write_text(html, encoding="utf-8")

        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(html.encode("utf-8")),
        )

    # ----- Audio (MP3 / script fallback) -----

    def assemble_audio(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """오디오 스크립트 HTML (미니멀 디자인 시스템)."""
        # MP3 에셋이 있으면 합성 시도
        if assets:
            mp3_assets = [a for a in assets if a.asset_type == "audio"]
            if mp3_assets and all(os.path.exists(a.file_path) for a in mp3_assets):
                try:
                    return self._merge_audio(mp3_assets, topic)
                except Exception as e:
                    logger.warning(f"오디오 합성 실패: {e}")

        # 스크립트 텍스트 fallback
        sections = content.sections if hasattr(content, "sections") else ()
        
        script_html = ""
        for sec in sections:
            script_html += f"""
            <div style="margin-bottom:30px; display:flex; gap:20px;">
                <div style="width:100px; flex-shrink:0; font-weight:800; color:#0052A2; text-align:right; font-size:14px; padding-top:5px;">{sec.speaker}</div>
                <div style="background:#F8F9FA; padding:20px; border-radius:4px; flex:1; font-size:16px; line-height:1.7; color:#333;">{sec.text.replace(chr(10), '<br>')}</div>
            </div>
            """

        html = f"""<!DOCTYPE html>
<html lang="ko">
<head><meta charset="utf-8">
<style>
    body{{font-family:'Pretendard','Apple SD Gothic Neo',sans-serif; background:#FFFFFF; color:#1A1A2E; margin:0; padding:60px; line-height:1.6;}}
    header{{border-bottom:2px solid #EEE; margin-bottom:50px; padding-bottom:30px;}}
    h1{{color:#0052A2; font-size:36px; margin:0 0 10px; font-weight:800;}}
    footer{{margin-top:100px; color:#AAA; font-size:12px; letter-spacing:1px; border-top:1px solid #EEE; padding-top:20px;}}
</style>
</head>
<body>
    <header>
        <h1>{topic}</h1>
        <div style="color:#666; font-size:16px;">오디오 교육 대본</div>
    </header>
    <div style="max-width:800px;">{script_html}</div>
    <footer>SKMS CONTENT STUDIO</footer>
</body></html>"""

        out_path = self._output_path("audio", topic, "html")
        out_path.write_text(html, encoding="utf-8")

        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(html.encode("utf-8")),
        )

    def _merge_audio(
        self, mp3_assets: list[GeneratedAsset], topic: str
    ) -> GeneratedFile:
        """MP3 파일 합성 (pydub 사용)."""
        from pydub import AudioSegment

        combined = AudioSegment.empty()
        for asset in mp3_assets:
            segment = AudioSegment.from_mp3(asset.file_path)
            combined += segment

        out_path = self._output_path("audio", topic, "mp3")
        combined.export(str(out_path), format="mp3")
        size = out_path.stat().st_size

        return GeneratedFile(
            file_type="mp3",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=size,
        )

    # ----- Video Studio (MP4) -----

    def assemble_video(
        self,
        content: Any,
        assets: list[GeneratedAsset],
        topic: str = "",
    ) -> GeneratedFile:
        """슬라이드 이미지와 오디오를 합성하여 MP4 영상 강의를 생성한다.

        PR-065: Video Studio Engine 기초 구현.
        """
        import subprocess
        import tempfile

        # 1. 에셋 분류 (이미지 vs 오디오)
        image_assets = sorted([a for a in assets if a.asset_type == "image"], 
                             key=lambda x: dict(x.metadata or {}).get("slide_index", 0))
        audio_assets = sorted([a for a in assets if a.asset_type == "audio"],
                             key=lambda x: dict(x.metadata or {}).get("slide_index", 0))

        if not image_assets or not audio_assets:
            logger.warning("이미지 또는 오디오 에셋 부족으로 영상 생성을 스킵합니다.")
            return None

        out_path = self._output_path("lectures", f"{topic}-video", "mp4")
        
        # 2. 임시 작업 디렉토리 생성
        with tempfile.TemporaryDirectory() as tmp_dir:
            # ffmpeg용 input 리스트 파일 생성
            concat_file = Path(tmp_dir) / "inputs.txt"
            
            with open(concat_file, "w", encoding="utf-8") as f:
                for img in image_assets:
                    f.write(f"file '{os.path.abspath(img.file_path)}'\n")
                    f.write("duration 5\n") # Day 2에서 동적 계산 로직 추가 예정
                # ffmpeg bug 방지 위해 마지막 파일 한 번 더 기재
                if image_assets:
                    f.write(f"file '{os.path.abspath(image_assets[-1].file_path)}'\n")

            # 3. FFmpeg 실행 (이미지 시퀀스 + 오디오 합성)
            try:
                # [참고] audio_assets[0]은 전체 합성된 오디오이거나 첫 번째 오디오일 수 있음
                # Day 2에서 오디오 합본 자동 생성 로직과 연동 예정
                cmd = [
                    "ffmpeg", "-y",
                    "-f", "concat", "-safe", "0", "-i", str(concat_file),
                    "-i", os.path.abspath(audio_assets[0].file_path),
                    "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", "25",
                    "-shortest",
                    str(out_path)
                ]
                subprocess.run(cmd, check=True, capture_output=True)
            except Exception as e:
                logger.error(f"영상 합성 실패: {e}")
                return None

        return GeneratedFile(
            file_type="mp4",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=out_path.stat().st_size if out_path.exists() else 0,
        )

    # ----- Visualization (SVG/PNG passthrough) -----

    def assemble_visualization(
        self,
        content: Any,
        assets: list[GeneratedAsset] | None = None,
        topic: str = "",
    ) -> GeneratedFile:
        """시각화 에셋 파일 복사."""
        if assets:
            chart_assets = [a for a in assets if a.asset_type == "chart"]
            if chart_assets:
                import shutil

                asset = chart_assets[0]
                ext = "svg" if asset.file_path.endswith(".svg") else "png"
                out_path = self._output_path("visualization", topic, ext)
                shutil.copy2(asset.file_path, str(out_path))
                return GeneratedFile(
                    file_type=ext,
                    file_path=str(out_path),
                    file_name=out_path.name,
                    size_bytes=out_path.stat().st_size,
                )

        # No asset → HTML timeline fallback from VisualizationPlan
        out_path = self._output_path("visualization", topic, "html")
        html = _build_visualization_html(content, topic)
        out_path.write_text(html, encoding="utf-8")
        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(html.encode("utf-8")),
        )

    # ----- Quiz (HTML) -----

    def assemble_quiz(
        self,
        content: Any,
        topic: str = "",
    ) -> GeneratedFile:
        """퀴즈 HTML 생성 — QuizPlan에서 인터랙티브 HTML."""
        out_path = self._output_path("quiz", topic, "html")
        html = _build_quiz_html(content, topic)
        out_path.write_text(html, encoding="utf-8")
        return GeneratedFile(
            file_type="html",
            file_path=str(out_path),
            file_name=out_path.name,
            size_bytes=len(html.encode("utf-8")),
        )


# ---------------------------------------------------------------------------
# HTML builders for quiz and visualization
# ---------------------------------------------------------------------------


def _build_quiz_html(plan: Any, topic: str) -> str:
    """QuizPlan → 인터랙티브 퀴즈 HTML."""
    questions = getattr(plan, "questions", ())
    if not questions:
        return f"<html><body><h1>{topic}</h1><p>퀴즈 문항이 없습니다.</p></body></html>"

    import html as html_mod

    question_blocks = []
    for q in questions:
        choices_html = ""
        for ci, choice in enumerate(q.choices):
            escaped = html_mod.escape(choice)
            choices_html += (
                f'<button class="choice" data-q="{q.index}" data-c="{ci}" '
                f'onclick="checkAnswer(this,{q.index},{q.correct_answer})">'
                f'<span class="label">{chr(65 + ci)}</span> {escaped}</button>\n'
            )
        explanation_escaped = html_mod.escape(q.explanation)
        source = html_mod.escape(getattr(q, "source_quote", ""))
        source_html = f'<div class="source">출처: {source}</div>' if source else ""
        question_blocks.append(
            f'<div class="question" id="q{q.index}">'
            f'<div class="q-num">Q{q.index}</div>'
            f'<div class="q-text">{html_mod.escape(q.question_text)}</div>'
            f'<div class="choices">{choices_html}</div>'
            f'<div class="explanation" id="exp{q.index}" style="display:none;">'
            f"{explanation_escaped}{source_html}</div>"
            f"</div>"
        )

    questions_joined = "\n".join(question_blocks)
    title_escaped = html_mod.escape(topic)
    n = len(questions)
    return (
        "<!DOCTYPE html>\n"
        '<html lang="ko">\n'
        '<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">\n'
        f"<title>{title_escaped}</title>\n"
        "<style>\n"
        "body{font-family:'Pretendard','Apple SD Gothic Neo',sans-serif;background:#FFFFFF;margin:0;padding:40px;color:#1A1A2E;line-height:1.6}\n"
        "h1{text-align:left;color:#0052A2;margin:0 0 40px;font-size:32px;font-weight:800;border-bottom:3px solid #0052A2;padding-bottom:15px}\n"
        ".question{background:white;border-radius:0;padding:30px 0;margin-bottom:30px;border-bottom:1px solid #EEE}\n"
        ".q-num{color:#0052A2;font-weight:800;font-size:16px;margin-bottom:10px;letter-spacing:1px}\n"
        ".q-text{font-size:20px;font-weight:600;margin-bottom:20px}\n"
        ".choices{display:flex;flex-direction:column;gap:12px}\n"
        ".choice{border:1px solid #DDD;border-radius:4px;padding:15px 20px;background:white;cursor:pointer;"
        "text-align:left;font-size:16px;transition:all 0.2s;color:#444}\n"
        ".choice:hover{border-color:#0052A2;background:#F0F7FF;color:#0052A2}\n"
        ".choice .label{display:inline-block;width:24px;height:24px;border-radius:2px;background:#F0F0F0;"
        "text-align:center;line-height:24px;font-weight:700;margin-right:12px;font-size:13px}\n"
        ".choice.correct{border-color:#28a745;background:#f8fff9;color:#28a745}"
        ".choice.correct .label{background:#28a745;color:white}\n"
        ".choice.wrong{border-color:#dc3545;background:#fff8f8;color:#dc3545}"
        ".choice.wrong .label{background:#dc3545;color:white}\n"
        ".choice.disabled{pointer-events:none;opacity:0.7}\n"
        ".explanation{margin-top:20px;padding:20px;background:#F8F9FA;border-left:5px solid #0052A2;"
        "font-size:15px;color:#555}\n"
        ".source{margin-top:10px;font-size:13px;color:#0052A2;font-weight:600}\n"
        ".score{text-align:center;font-size:24px;font-weight:800;color:#0052A2;margin-top:50px;padding:30px;background:#F0F7FF;border-radius:8px}\n"
        "footer{text-align:left;margin-top:60px;color:#AAA;font-size:12px;letter-spacing:1px}\n"
        "</style>\n"
        "</head>\n<body>\n"
        f"<h1>{title_escaped}</h1>\n"
        f"{questions_joined}\n"
        '<div class="score" id="score"></div>\n'
        "<footer>SKMS Content Studio</footer>\n"
        "<script>\n"
        f"let answered=0,correct=0,total={n};\n"
        "function checkAnswer(btn,qIdx,correctIdx){\n"
        "  let btns=document.querySelectorAll('.choice[data-q=\"'+qIdx+'\"]');\n"
        "  btns.forEach(b=>{b.classList.add('disabled');"
        "if(parseInt(b.dataset.c)===correctIdx)b.classList.add('correct')});\n"
        "  if(parseInt(btn.dataset.c)!==correctIdx)btn.classList.add('wrong');\n"
        "  else correct++;\n"
        "  answered++;\n"
        "  document.getElementById('exp'+qIdx).style.display='block';\n"
        "  if(answered===total)document.getElementById('score').textContent="
        "'결과: '+total+'문제 중 '+correct+'개 정답 ('+Math.round(correct/total*100)+'%)';\n"
        "}\n"
        "</script>\n"
        "</body></html>"
    )


def _build_visualization_html(plan: Any, topic: str) -> str:
    """VisualizationPlan → HTML 타임라인."""
    import html as html_mod

    title = getattr(plan, "title", topic) or topic
    description = getattr(plan, "data_description", "") or ""

    editions = [
        ("1979", "초판", "SK경영체계의 출발점"),
        ("1981", "1차", "경영 원칙 체계화"),
        ("1988", "5차", "경영 요소 분류 재편"),
        ("1989", "6차", "SUPEX 추구 도입"),
        ("1990", "7차", "관리 체계 정교화"),
        ("1995", "8차", "글로벌 경영 반영"),
        ("1997", "9차", "SUPEX 추구법 완성"),
        ("1998", "10차", "구조조정기 반영"),
        ("2004", "11차", "지배구조 변화 반영"),
        ("2008", "12차", "사회적 책임 강화"),
        ("2016", "13차", "SUPEX Company 개념"),
        ("2020", "14차", "VWBE 문화 · 구성원 행복"),
    ]

    items = ""
    for year, name, desc in editions:
        items += (
            f'<div class="tl-item">'
            f'<div class="tl-year">{year}</div>'
            f'<div class="tl-dot"></div>'
            f'<div class="tl-card">'
            f'<div class="tl-name">{html_mod.escape(name)} 개정판</div>'
            f'<div class="tl-desc">{html_mod.escape(desc)}</div>'
            f"</div></div>\n"
        )

    title_escaped = html_mod.escape(title)
    desc_escaped = html_mod.escape(description)
    return (
        "<!DOCTYPE html>\n"
        '<html lang="ko">\n'
        '<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">\n'
        f"<title>{title_escaped}</title>\n"
        "<style>\n"
        "body{font-family:'Pretendard','Apple SD Gothic Neo',sans-serif;background:#FFFFFF;margin:0;padding:40px;color:#1A1A2E}\n"
        "h1{text-align:left;color:#0052A2;margin:0 0 10px;font-size:32px;font-weight:800}\n"
        ".subtitle{text-align:left;color:#666;margin-bottom:50px;font-size:16px;border-left:4px solid #0052A2;padding-left:15px}\n"
        ".timeline{position:relative;max-width:800px;margin:0;padding:20px 0}\n"
        ".timeline::before{content:'';position:absolute;left:100px;top:0;bottom:0;width:2px;"
        "background:#EEE}\n"
        ".tl-item{display:flex;align-items:flex-start;margin-bottom:40px;position:relative}\n"
        ".tl-year{width:80px;text-align:right;font-weight:800;color:#0052A2;font-size:18px;flex-shrink:0;padding-top:10px}\n"
        ".tl-dot{width:12px;height:12px;border-radius:50%;background:#FFFFFF;border:3px solid #0052A2;"
        "margin:15px 20px 0;flex-shrink:0;z-index:1}\n"
        ".tl-card{background:#F8F9FA;border-radius:0;padding:20px 25px;flex:1;border-left:2px solid #0052A2}\n"
        ".tl-name{font-weight:700;font-size:18px;color:#1A1A2E}\n"
        ".tl-desc{font-size:15px;color:#444;margin-top:8px;line-height:1.6}\n"
        "footer{text-align:left;margin-top:80px;color:#AAA;font-size:12px;letter-spacing:1px}\n"
        "</style>\n"
        "</head>\n<body>\n"
        f"<h1>{title_escaped}</h1>\n"
        f'<div class="subtitle">{desc_escaped}</div>\n'
        f'<div class="timeline">\n{items}</div>\n'
        "<footer>SKMS Content Studio</footer>\n"
        "</body></html>"
    )
